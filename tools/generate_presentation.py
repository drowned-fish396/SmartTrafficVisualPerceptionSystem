from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = "docs/需求分析与系统设计V1.0汇报.pptx"


def set_run(run, size=18, bold=False):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.25), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_run(p.runs[0], 26, True)


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.25), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"需求分析报告 V1.0 / 系统设计方案 V1.0    {page}/13"
    set_run(p.runs[0], 9, False)
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, bullets, x=0.75, y=1.15, w=12, h=5.4, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.runs[0], size, False)


def add_section_label(slide, text, x, y, w, h, size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], size, True)
    return shape


def add_table_like(slide, rows, x=0.7, y=1.2, w=11.9, row_h=0.55, size=13):
    col_count = len(rows[0])
    col_w = w / col_count
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + c * col_w),
                Inches(y + r * row_h),
                Inches(col_w),
                Inches(row_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.line.width = Pt(0.7)
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell
            p.alignment = PP_ALIGN.CENTER
            set_run(p.runs[0], size, r == 0)


def add_er_slide(slide):
    add_placeholder(
        slide,
        "此处插入：数据库 ER 图",
        [
            "建议放在页面左侧或上方：展示设备、视频源、分析任务、模型配置、检测结果、告警与统计数据之间的关系。",
            "核心关系：设备关联视频源和分析任务；视频源进入分析任务；模型配置被分析任务引用；分析任务产生检测结果和业务事件。",
        ],
        x=0.7,
        y=1.05,
        w=5.4,
        h=5.25,
    )
    rows = [
        ["表名", "用途"],
        ["devices", "存储边端设备信息"],
        ["video_sources", "存储实时流和上传视频信息"],
        ["analysis_tasks", "存储视频分析任务信息"],
        ["model_configs", "存储模型配置和运行参数"],
        ["whitelist_vehicles", "存储白名单车辆库"],
        ["detection_results", "存储通用检测结果"],
        ["plate_records", "存储车牌识别和白名单比对记录"],
        ["congestion_stats", "存储拥堵统计和热力图数据"],
        ["parking_events", "存储禁停区域车辆停留和告警事件"],
        ["anomaly_events", "存储道路异常检测事件"],
        ["system_metrics", "存储系统资源和运行状态"],
        ["operation_logs", "存储关键操作和异常日志"],
    ]
    add_table_like(slide, rows, x=6.35, y=1.05, w=6.2, row_h=0.4, size=8.8)


def add_placeholder(slide, title, details, x=0.8, y=1.25, w=11.7, h=4.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], 22, True)
    for detail in details:
        p = tf.add_paragraph()
        p.text = detail
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)
        set_run(p.runs[0], 15, False)


def add_slide(prs, page, title, bullets=None, table=None, placeholder=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_title(slide, title)
    if bullets:
        add_bullets(slide, bullets)
    if table:
        add_table_like(slide, table)
    if placeholder:
        add_placeholder(slide, placeholder[0], placeholder[1])
    add_footer(slide, page)
    return slide


def add_er_table_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_title(slide, "ER 图与核心数据表设计")
    add_er_slide(slide)
    add_footer(slide, page)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_slide(
        prs,
        1,
        "云边端协同的智慧交通视觉感知系统",
        [
            "汇报范围：需求分析报告 V1.0 与系统设计方案 V1.0 的核心内容。",
            "展示重点：关键需求决策、需求分析用例图、功能架构图、技术架构图、ER 图。",
            "系统定位：面向 704 智慧交通沙盘，构建手机边端、云端分析服务、前端展示端协同运行的平台。",
            "汇报目标：说明系统为什么这样设计，以及这些设计如何支撑课程演示闭环。",
        ],
    )

    add_slide(
        prs,
        2,
        "项目目标与演示闭环",
        [
            "项目目标：完成从手机视频采集、视频传输、云端 AI 分析、前端展示到历史查询的完整链路。",
            "核心闭环：边端采集视频 -> 云端接收与分析 -> 前端实时展示 -> 数据持久化与历史查询。",
            "演示场景：优先面向学院 704 智慧交通沙盘，也可兼容校园道路或校外道路画面。",
            "阶段原则：V1.0 优先保证可演示闭环，再扩展多设备、多模型、多场景和更丰富的管理能力。",
            "系统边界：不要求真实控制闸机，不承诺生产级复杂道路准确率，允许视频段上传和模拟结果兜底。",
        ],
    )

    add_slide(
        prs,
        3,
        "需求分析用例图",
        placeholder=(
            "此处插入：需求分析用例图",
            [
                "参与者：边端设备使用者、平台管理员、算法/模型配置人员、监控查看人员。",
                "边端用例：采集交通视频、推送实时视频流、上传短视频段、设备注册、上报心跳。",
                "管理用例：查看设备状态、查看系统状态、配置模型参数、配置检测区域、配置禁停区域。",
                "监控用例：查看实时监控、车牌识别、拥堵热力图、禁停告警、道路异常告警、历史数据。",
            ],
        ),
    )

    add_slide(
        prs,
        4,
        "关键需求决策总览",
        table=[
            ["关键问题", "核心决策", "解决的主要风险"],
            ["数据传输方式", "RTSP/RTMP + REST API + WebSocket + HTTP 上传", "避免单一协议承载所有数据"],
            ["实时性与连续性", "实时流优先，抽帧分析，视频段与模拟结果兜底", "网络、算力或模型不稳定导致演示中断"],
            ["车牌数字识别", "YOLO 定位车牌，PaddleOCR 识别完整车牌并提取数字", "车牌模糊、误识别和白名单比对不可信"],
            ["热力图形式", "分区密度热力图 + 时间窗口统计", "热力图只有装饰效果、缺少统计依据"],
            ["跟踪与异常检测", "YOLO + ByteTrack；FastFlow + 规则检测兜底", "禁停无法持续计时、未知异常类别难以穷举"],
        ],
    )

    add_slide(
        prs,
        5,
        "决策 1：数据传输方式权衡",
        table=[
            ["数据类型", "选择方式", "原因"],
            ["实时视频流", "RTSP/RTMP", "视频数据体量大、连续性强，适合边端向云端推送或云端拉流"],
            ["设备注册、心跳、配置、查询", "REST API", "请求响应清晰，便于接口测试、权限控制和文档维护"],
            ["检测框、车牌、热力图、告警、状态", "WebSocket", "云端主动推送，减少前端轮询开销，保证监控页面及时刷新"],
            ["实时流异常输入", "HTTP 视频上传", "网络或推流环境不稳定时，保证分析和展示流程继续运行"],
            ["结论", "按数据特点拆分通道", "视频、管理、实时结果各走最合适的传输方式"],
        ],
    )

    add_slide(
        prs,
        6,
        "决策 2：实时性与连续性权衡",
        [
            "主方案：边端手机优先推送实时视频流，云端连续拉流并按配置频率抽取视频帧。",
            "实时性目标：端到端延迟目标为 1-2 秒，理想分析帧率目标为 15-30 FPS。",
            "降级策略：硬件算力不足时降低抽帧频率、缩小输入分辨率、切换轻量模型，控制在 5-15 FPS 范围内保持连续展示。",
            "连续性兜底：实时流中断或推流条件不满足时，使用短视频段上传进入同一套分析任务流程。",
            "演示保障：模型文件、GPU 或识别效果不稳定时，使用结构一致的模拟结果，避免前端和数据库逻辑分裂。",
        ],
    )

    add_slide(
        prs,
        7,
        "决策 3：车牌识别如何识别数字",
        [
            "技术路线：YOLO 车牌定位 + PaddleOCR 字符识别 + 车牌规则校验。",
            "处理流程：视频帧 -> 车辆/车牌检测 -> 车牌区域裁剪 -> 图像增强与校正 -> PaddleOCR 输出完整车牌字符串。",
            "数字识别方式：不单独训练数字模型，而是从 PaddleOCR 的完整车牌结果中提取数字字符，同时保留完整车牌号。",
            "可信度控制：结合长度校验、字符类别校验和 OCR 置信度阈值，提高白名单比对结果可信度。",
            "异常处理：车牌过小、模糊、遮挡或置信度不足时标记为“需人工确认”，演示模式可使用预设样例结果。",
        ],
    )

    add_slide(
        prs,
        8,
        "决策 4：热力图形式",
        [
            "热力图方案：采用分区密度热力图，不单独训练热力图模型，也不只做装饰性渐变效果。",
            "区域划分：将 704 沙盘道路划分为入口区、路口区、闸机等待区、转弯区、禁停区域周边等检测区域。",
            "统计依据：根据 YOLO 车辆检测框中心点所属区域统计车辆数量，并用“区域车辆数 / 区域容量”计算密度等级。",
            "稳定策略：结合时间窗口统计，降低单帧误检或漏检造成的热力图抖动。",
            "前端展示：视频底图层 + 半透明热力覆盖层 + 区域名称、车辆数、趋势箭头和持续时间标注层。",
        ],
    )

    add_slide(
        prs,
        9,
        "决策 5：车辆跟踪模型与异常检测模型选择",
        [
            "禁停跟踪路线：YOLO 车辆检测 + ByteTrack 跟踪 + 禁停区域规则。",
            "选择理由：禁停判断需要识别同一车辆的持续 ID、进入时间、当前位置和停留时长，单帧检测无法完成持续计时。",
            "告警逻辑：车辆在禁停区域内停留超过预设阈值时，生成长时间停留告警，并展示轨迹、跟踪 ID 和告警级别。",
            "异常检测路线：FastFlow 异常检测模型 + 规则检测兜底。",
            "选择理由：道路异常物体可能未知，难以提前穷举类别，因此不把常规目标检测作为异常检测主方案。",
            "兜底方案：基于固定视角、道路区域掩膜、背景差分、静止区域变化、异常面积阈值和持续时间阈值判断疑似异常。",
        ],
    )

    add_slide(
        prs,
        10,
        "功能架构图",
        placeholder=(
            "此处插入：功能架构图",
            [
                "边端：视频采集、实时推流、视频上传、设备注册、心跳上报、状态提示。",
                "云端接入层：REST API、WebSocket、RTSP/RTMP 接入、HTTP 文件上传。",
                "云端业务层：设备管理、视频源管理、分析任务、模型管理、系统监控、告警管理、历史查询。",
                "云端算法层：车辆检测、车牌识别、拥堵分析、禁停跟踪、道路异常检测、模拟结果兜底。",
                "前端：实时监控大屏、闸机监控、拥堵热力图、禁停监控、异常监控、设备/模型/系统/历史页面。",
            ],
        ),
    )

    add_slide(
        prs,
        11,
        "技术架构图",
        placeholder=(
            "此处插入：技术架构图",
            [
                "边端采集层：手机摄像头、推流工具、HTTP 上传。",
                "云端接入层：FastAPI、Uvicorn、WebSocket、RTSP/RTMP、HTTP Multipart。",
                "云端业务层：FastAPI 路由、Pydantic、psutil。",
                "云端算法层：YOLOv8n、PaddleOCR、ByteTrack、FastFlow、OpenCV、FFmpeg。",
                "云端数据层：SQLAlchemy、SQLite、文件系统；前端展示层：Vue 3、TypeScript、Vite、ECharts、Axios、WebSocket。",
            ],
        ),
    )

    add_er_table_slide(prs, 12)

    add_slide(
        prs,
        13,
        "总结：决策驱动系统设计",
        [
            "需求层面：系统围绕 704 智慧交通沙盘，覆盖车牌识别、拥堵热力图、禁停跟踪、道路异常检测和平台管理能力。",
            "架构层面：采用云边端协同的分层模块化架构，明确边端采集、云端分析、前端展示的职责边界。",
            "技术层面：通过 RTSP/RTMP、REST API、WebSocket、HTTP 上传组合，支撑视频流、业务管理和实时结果推送。",
            "算法层面：使用 YOLO、PaddleOCR、ByteTrack、FastFlow 和规则兜底，兼顾识别能力、实时性、可解释性和演示稳定性。",
            "验收层面：即使网络、GPU 或模型不稳定，也能通过短视频上传和模拟结果保证课程演示闭环。",
        ],
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
